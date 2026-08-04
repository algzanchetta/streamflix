from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def criar_powerpoint():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Função auxiliar para adicionar slide com título e conteúdo
    def adicionar_slide(titulo, conteudo, notas=""):
        slide_layout = prs.slide_layouts[1]  # Título e conteúdo
        slide = prs.slides.add_slide(slide_layout)
        
        # Título
        titulo_slide = slide.shapes.title
        titulo_slide.text = titulo
        titulo_slide.text_frame.paragraphs[0].font.size = Pt(44)
        titulo_slide.text_frame.paragraphs[0].font.bold = True
        titulo_slide.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Conteúdo
        corpo = slide.placeholders[1]
        tf = corpo.text_frame
        tf.text = conteudo
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(24)
            paragraph.space_after = Pt(12)
        
        # Notas do apresentador (aparecem no modo de apresentação)
        if notas:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notas
    
    # SLIDE 1: CAPA
    slide_layout = prs.slide_layouts[0]  # Slide de título
    slide = prs.slides.add_slide(slide_layout)
    titulo = slide.shapes.title
    titulo.text = "Sistema Web de Streaming\ncom Python e Flask"
    titulo.text_frame.paragraphs[0].font.size = Pt(54)
    titulo.text_frame.paragraphs[0].font.bold = True
    titulo.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    subtítulo = slide.placeholders[1]
    subtítulo.text = "Aula 00 - Apresentação do Projeto + Preparação do Banco de Dados"
    subtítulo.text_frame.paragraphs[0].font.size = Pt(28)
    
    # SLIDE 2: O QUE VAMOS CONSTRUIR?
    adicionar_slide(
        "O Que Vamos Construir?",
        "✅ Sistema web profissional para gestão de streaming (StreamFlix)\n\n"
        "✅ Área do Cliente: login, planos, pagamentos, suporte\n\n"
        "✅ Área Administrativa: dashboard, métricas, CRUD completo\n\n"
        "✅ API RESTful, deploy, testes e boas práticas\n\n"
        "✅ Projeto real e completo para o seu portfólio",
        "Não vamos fazer apenas um hello world. Vamos construir um sistema completo, com autenticação, painel administrativo e relatórios. Tudo baseado em uma empresa real de streaming."
    )
    
    # SLIDE 3: O BANCO DE DADOS ATUAL
    adicionar_slide(
        "O Banco de Dados Atual",
        " Tabelas existentes:\n"
        "   • clientes (id, nome, email, idade, cidade)\n"
        "   • assinaturas (id, id_cliente, plano, valor, status)\n"
        "   • pagamentos (id, id_assinatura, data, valor, método)\n\n"
        "❓ O que falta para virar um sistema web?\n"
        "   • Login seguro\n"
        "   • Logs de auditoria\n"
        "   • Tickets de suporte\n"
        "   • Configurações dinâmicas\n"
        "   • Notificações",
        "Atualmente temos um banco funcional para análise de dados, mas para virar um sistema web, precisamos adicionar camadas de segurança e gestão. É isso que faremos agora."
    )
    
    # SLIDE 4: JORNADA DO CURSO
    adicionar_slide(
        "Jornada do Curso (Ementa)",
        "🏗️ Módulo 1: Fundações (Setup, Models, Auth) - 5 aulas\n\n"
        "👤 Módulo 2: Área do Cliente (Dashboard, Planos, Suporte) - 7 aulas\n\n"
        "⚙️ Módulo 3: Área Admin (Métricas, CRUD, Relatórios) - 8 aulas\n\n"
        "🚀 Módulo 4: Avançado (API, Pagamentos, Cache, Deploy) - 8 aulas\n\n"
        "🎁 Módulo 5: Bônus (Dark Mode, PWA) - 2 aulas\n\n"
        "📚 Total: 30+ aulas práticas",
        "São mais de 30 aulas práticas. Cada módulo constrói sobre o anterior. No final, você terá um sistema completo pronto para mostrar em entrevistas."
    )
    
    # SLIDE 5: PREPARAÇÃO DO BANCO
    adicionar_slide(
        "Preparação e Migração do Banco",
        " O que faremos agora:\n\n"
        "1. Executar script de migração SQL\n\n"
        "2. Recriar tabela clientes com novos campos\n"
        "   (senha_hash, data_cadastro, ultimo_acesso, ativo)\n\n"
        "3. Criar 7 tabelas novas\n"
        "   (usuários, logs, suporte, configurações, etc.)\n\n"
        "4. Adicionar índices e views para relatórios\n\n"
        "⚠️ Importante: Faça backup do arquivo dbase.db antes!",
        "O SQLite não permite adicionar certas colunas via ALTER TABLE. Por isso, o script recria a tabela clientes de forma segura, preservando todos os dados."
    )
    
    # SLIDE 6: PRÓXIMOS PASSOS
    adicionar_slide(
        "Próximos Passos",
        "✅ Para você fazer agora:\n"
        "   • Baixar o script SQL e a ementa (links na descrição)\n"
        "   • Executar a migração no seu banco\n"
        "   • Deixar o like, se inscrever\n"
        "   • Comentar: 'Migração concluída!'\n\n"
        "🎯 Próxima aula (01):\n"
        "   • Configuração do ambiente virtual Python\n"
        "   • Estrutura de pastas profissional\n"
        "   • Primeiro setup do Flask + SQLAlchemy\n"
        "   • 'Hello World' estruturado",
        "Na aula 1 já começamos a codar de verdade. Tudo estará pronto para receber os models, rotas e templates. Não perca! Ative o sininho."
    )
    
    # SLIDE 7: ENCERRAMENTO
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    titulo = slide.shapes.title
    titulo.text = "Obrigado!"
    titulo.text_frame.paragraphs[0].font.size = Pt(60)
    titulo.text_frame.paragraphs[0].font.bold = True
    titulo.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    subtítulo = slide.placeholders[1]
    subtítulo.text = "Nos vemos na Aula 01! 🚀\n\nMaterial de apoio disponível na descrição"
    subtítulo.text_frame.paragraphs[0].font.size = Pt(32)

    # Salvar o arquivo
    prs.save('Aula_00_Apresentacao.pptx')
    print("✅ PowerPoint 'Aula_00_Apresentacao.pptx' gerado com sucesso!")

if __name__ == '__main__':
    criar_powerpoint()